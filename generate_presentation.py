#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Génère la présentation PPTX professionnelle du projet ENT EST Salé."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BLEU_FONCE = RGBColor(0x0F, 0x29, 0x4D)
BLEU_MED   = RGBColor(0x1A, 0x56, 0xDB)
BLEU_CLAIR = RGBColor(0xDB, 0xEA, 0xFE)
VERT       = RGBColor(0x05, 0x96, 0x69)
VERT_CLAIR = RGBColor(0xD1, 0xFA, 0xE5)
ORANGE     = RGBColor(0xD9, 0x77, 0x06)
ORANGE_CL  = RGBColor(0xFE, 0xF3, 0xC7)
ROUGE      = RGBColor(0xDC, 0x26, 0x26)
GRIS       = RGBColor(0x37, 0x41, 0x51)
GRIS_CLAIR = RGBColor(0xF1, 0xF5, 0xF9)
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
FOND       = RGBColor(0xF8, 0xFA, 0xFF)

W, H = Inches(13.33), Inches(7.5)   # Widescreen 16:9


# ── Helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_shape_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, radius=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    set_shape_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 size=18, bold=False, italic=False,
                 color: RGBColor = GRIS,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para_to_tf(tf, text, size=14, bold=False, color: RGBColor = GRIS,
                   align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, color=BLEU_FONCE, height=Inches(0.08)):
    add_rect(slide, 0, 0, W, height, color)


def add_left_accent(slide, color=BLEU_MED, width=Inches(0.06)):
    add_rect(slide, 0, 0, width, H, color)


def card(slide, left, top, width, height, bg=BLANC, border_color=BLEU_MED, title="", title_size=13, body="", body_size=11):
    r = add_rect(slide, left, top, width, height, bg)
    # left border accent
    add_rect(slide, left, top, Inches(0.05), height, border_color)
    if title:
        add_text_box(slide, title, left + Inches(0.15), top + Inches(0.1),
                     width - Inches(0.2), Inches(0.35),
                     size=title_size, bold=True, color=border_color)
    if body:
        add_text_box(slide, body, left + Inches(0.15), top + Inches(0.45),
                     width - Inches(0.2), height - Inches(0.55),
                     size=body_size, color=GRIS)
    return r


def badge(slide, text, left, top, w=Inches(1.4), h=Inches(0.35), bg=BLEU_MED, txt_color=BLANC, size=10):
    add_rect(slide, left, top, w, h, bg)
    add_text_box(slide, text, left, top, w, h, size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
def build_presentation():
    prs = new_prs()

    # ────────────────────────────────────────────────────────────── SLIDE 1: TITRE
    sl = blank_slide(prs)
    slide_bg(sl, BLEU_FONCE)

    # Gradient-like strip right
    add_rect(sl, W * 0.62, 0, W * 0.38, H, BLEU_MED)
    add_rect(sl, W * 0.60, 0, Inches(0.04), H, BLANC)

    # University
    add_text_box(sl, "Université Mohammed V — Rabat", Inches(0.5), Inches(0.4), Inches(7.5), Inches(0.5),
                 size=14, color=BLEU_CLAIR, align=PP_ALIGN.LEFT)
    add_text_box(sl, "École Supérieure de Technologie de Salé (EST Salé)", Inches(0.5), Inches(0.8), Inches(7.5), Inches(0.5),
                 size=12, color=BLANC, align=PP_ALIGN.LEFT)
    add_text_box(sl, "Filière : Ingénierie des Applications Web et Mobile  |  Module : DevOps et Cloud",
                 Inches(0.5), Inches(1.15), Inches(7.5), Inches(0.4), size=10, color=BLEU_CLAIR, align=PP_ALIGN.LEFT)

    # Title
    add_text_box(sl, "ENT-EST-Salé", Inches(0.5), Inches(2.0), Inches(7.5), Inches(1.2),
                 size=54, bold=True, color=BLANC)
    add_text_box(sl, "Plateforme Numérique de Travail", Inches(0.5), Inches(3.1), Inches(7.5), Inches(0.7),
                 size=22, color=BLEU_CLAIR)
    add_text_box(sl, "Architecture Microservices · DevOps · Cloud", Inches(0.5), Inches(3.75), Inches(7.5), Inches(0.5),
                 size=14, italic=True, color=BLEU_CLAIR)

    # Right panel info
    add_text_box(sl, "GROUPE", W * 0.64, Inches(0.5), Inches(4.5), Inches(0.4), size=11, bold=True, color=BLEU_CLAIR)
    membres = ["Membre 1", "Membre 2", "Membre 3", "Membre 4", "Membre 5"]
    for i, m in enumerate(membres):
        add_text_box(sl, f"  {m}", W * 0.64, Inches(0.9 + i * 0.4), Inches(4.5), Inches(0.38), size=13, color=BLANC)

    add_text_box(sl, "ENCADRANT", W * 0.64, Inches(3.2), Inches(4.5), Inches(0.4), size=11, bold=True, color=BLEU_CLAIR)
    add_text_box(sl, "  —", W * 0.64, Inches(3.6), Inches(4.5), Inches(0.38), size=13, color=BLANC)

    add_text_box(sl, "Année universitaire 2025–2026", W * 0.64, Inches(4.5), Inches(4.5), Inches(0.4), size=11, color=BLEU_CLAIR)
    add_text_box(sl, "Avril 2026", W * 0.64, Inches(4.85), Inches(4.5), Inches(0.4), size=13, bold=True, color=BLANC)

    # Bottom bar
    add_rect(sl, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x09, 0x1B, 0x36))
    add_text_box(sl, "DevOps et Cloud  ·  EST Salé  ·  Université Mohammed V", 0, H - Inches(0.45), W, Inches(0.4),
                 size=10, color=BLEU_CLAIR, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 2: PLAN
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_FONCE, Inches(0.08))
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "Plan de la Présentation", Inches(0.3), Inches(0.2), Inches(12), Inches(0.6),
                 size=26, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.85), Inches(1.5), Inches(0.04), BLEU_MED)

    plan_items = [
        ("01", "Contexte et Objectifs", BLEU_FONCE),
        ("02", "Architecture Générale", BLEU_MED),
        ("03", "Infrastructure Technique", VERT),
        ("04", "Microservices — Vue Détaillée", ORANGE),
        ("05", "API Gateway & Sécurité", BLEU_MED),
        ("06", "Interface Utilisateur (Frontend)", VERT),
        ("07", "Observabilité & Monitoring", ORANGE),
        ("08", "Tests & Déploiement", ROUGE),
        ("09", "Conclusion & Perspectives", BLEU_FONCE),
    ]

    cols = 3
    bw = Inches(3.9); bh = Inches(0.8)
    bmargin_x = Inches(0.35); bmargin_y = Inches(1.0)
    gap_x = Inches(0.15); gap_y = Inches(0.12)

    for i, (num, title, color) in enumerate(plan_items):
        col = i % cols
        row = i // cols
        x = bmargin_x + col * (bw + gap_x)
        y = bmargin_y + row * (bh + gap_y)
        add_rect(sl, x, y, bw, bh, BLANC)
        add_rect(sl, x, y, Inches(0.06), bh, color)
        add_text_box(sl, num, x + Inches(0.1), y + Inches(0.08), Inches(0.5), Inches(0.4),
                     size=20, bold=True, color=color)
        add_text_box(sl, title, x + Inches(0.55), y + Inches(0.18), bw - Inches(0.6), Inches(0.5),
                     size=13, bold=True, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 3: CONTEXTE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_FONCE)
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "01  Contexte et Objectifs", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), BLEU_MED)

    # Left: context
    card(sl, Inches(0.3), Inches(1.0), Inches(5.8), Inches(2.1), bg=BLANC, border_color=BLEU_FONCE,
         title="Contexte Académique", title_size=13,
         body="Projet de fin de module DevOps et Cloud — EST Salé (Université Mohammed V, Rabat).\nFilière : Ingénierie des Applications Web et Mobile.\nGroupe de 5 étudiants + encadrant.", body_size=11)
    card(sl, Inches(0.3), Inches(3.2), Inches(5.8), Inches(2.0), bg=BLANC, border_color=ORANGE,
         title="Problématique", title_size=13,
         body="Les établissements d'enseignement supérieur manquent d'une plateforme numérique unifiée pour :\n• Gestion des cours et contenus pédagogiques\n• Devoirs, soumissions et notation\n• Communication et forum académique\n• Calendrier et notifications", body_size=11)

    # Right: tech objectives
    card(sl, Inches(6.4), Inches(1.0), Inches(6.6), Inches(5.7), bg=BLANC, border_color=VERT,
         title="Objectifs Techniques", title_size=13)
    objectives = [
        "Architecture Microservices découplée",
        "Authentification OAuth2/OIDC (Keycloak + JWT)",
        "Persistance NoSQL avec Apache Cassandra",
        "Messagerie asynchrone RabbitMQ",
        "Stockage objet S3 avec MinIO",
        "Interface React moderne (shadcn/ui)",
        "Chat temps-réel WebSocket",
        "Observabilité Prometheus + Grafana",
        "Déploiement Docker Compose automatisé",
        "Smoke-tests fonctionnels complets",
    ]
    for i, obj in enumerate(objectives):
        add_rect(sl, Inches(6.55), Inches(1.55 + i * 0.5), Inches(0.2), Inches(0.2), VERT)
        add_text_box(sl, obj, Inches(6.85), Inches(1.5 + i * 0.5), Inches(6.0), Inches(0.38),
                     size=11, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 4: ARCHITECTURE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_FONCE)
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "02  Architecture Générale — Vue d'Ensemble", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=22, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(2.0), Inches(0.04), BLEU_MED)

    layers_arch = [
        ("PRÉSENTATION", "React 18 + TypeScript + Vite  |  Tailwind CSS + shadcn/ui  |  WebSocket", BLEU_FONCE, BLEU_CLAIR, Inches(0.25)),
        ("API GATEWAY", "ms-api-gateway (port 8008)  |  JWT Guard  |  CORS  |  Retry + Timeout", BLEU_MED, RGBColor(0xDB,0xEA,0xFE), Inches(0.25)),
        ("MICROSERVICES", "auth-core · identity-admin · course-content · course-access · notification\ncalendar-schedule · forum-chat · exam-assignment · ai-assistant", VERT, VERT_CLAIR, Inches(0.25)),
        ("INFRASTRUCTURE", "Keycloak · Cassandra · RabbitMQ · MinIO · Redis · Mailpit · Prometheus · Grafana", ORANGE, ORANGE_CL, Inches(0.25)),
    ]

    for i, (label, desc, col, bg_col, left) in enumerate(layers_arch):
        y = Inches(1.05 + i * 1.38)
        h = Inches(1.22)
        add_rect(sl, left, y, W - Inches(0.5), h, bg_col)
        add_rect(sl, left, y, Inches(0.08), h, col)
        add_text_box(sl, label, left + Inches(0.18), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     size=12, bold=True, color=col)
        add_text_box(sl, desc, left + Inches(0.18), y + Inches(0.5), W - Inches(1.2), Inches(0.65),
                     size=11, color=GRIS)
        if i < 3:
            arrow_x = Inches(6.5); arrow_y = y + h
            add_text_box(sl, "↕", arrow_x, arrow_y, Inches(0.5), Inches(0.2), size=14, bold=True, color=GRIS, align=PP_ALIGN.CENTER)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 5: SCHEMA ARCHI DETAILLE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_FONCE)
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "02  Schéma d'Architecture — Flux de Communication", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=22, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(2.0), Inches(0.04), BLEU_MED)

    # Draw architecture diagram
    # Browser
    add_rect(sl, Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.7), BLEU_FONCE)
    add_text_box(sl, "🌐 Navigateur\nReact SPA", Inches(0.3), Inches(1.1), Inches(2.0), Inches(0.7),
                 size=10, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text_box(sl, "HTTP/REST\n+ WebSocket", Inches(2.35), Inches(1.3), Inches(1.0), Inches(0.4),
                 size=8, italic=True, color=BLEU_MED, align=PP_ALIGN.CENTER)
    add_text_box(sl, "──────►", Inches(2.35), Inches(1.35), Inches(1.0), Inches(0.3),
                 size=12, color=BLEU_MED, align=PP_ALIGN.CENTER)

    # Gateway
    add_rect(sl, Inches(3.4), Inches(1.0), Inches(2.2), Inches(0.9), BLEU_MED)
    add_text_box(sl, "API Gateway\nPort 8008", Inches(3.4), Inches(1.0), Inches(2.2), Inches(0.9),
                 size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    # Auth core below gateway
    add_text_box(sl, "validates JWT", Inches(3.4), Inches(1.95), Inches(2.2), Inches(0.3),
                 size=8, italic=True, color=GRIS, align=PP_ALIGN.CENTER)
    add_text_box(sl, "↓", Inches(4.45), Inches(2.2), Inches(0.3), Inches(0.3),
                 size=14, color=BLEU_MED, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.55), Inches(2.5), Inches(1.9), Inches(0.65), BLEU_CLAIR)
    add_text_box(sl, "ms-auth-core\n(JWKS Keycloak)", Inches(3.55), Inches(2.5), Inches(1.9), Inches(0.65),
                 size=9, color=BLEU_FONCE, align=PP_ALIGN.CENTER)

    # Services grid
    svcs = [
        ("identity-admin\n8013", BLEU_FONCE),
        ("course-content\n8011", BLEU_MED),
        ("course-access\n8012", BLEU_MED),
        ("notification\n8014", VERT),
        ("calendar\n8015", VERT),
        ("forum-chat\n8016", ORANGE),
        ("exam-assignment\n8017", ORANGE),
        ("ai-assistant\n8018", ROUGE),
    ]
    add_text_box(sl, "──── proxifie vers ────►", Inches(5.65), Inches(1.35), Inches(1.5), Inches(0.3),
                 size=8, italic=True, color=GRIS)

    for i, (name, col) in enumerate(svcs):
        col_i = i % 4; row_i = i // 4
        x = Inches(7.3) + col_i * Inches(1.47)
        y = Inches(1.05) + row_i * Inches(0.95)
        add_rect(sl, x, y, Inches(1.35), Inches(0.8), col)
        add_text_box(sl, name, x, y, Inches(1.35), Inches(0.8),
                     size=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    # Infra bottom row
    infra_svcs = [
        ("Keycloak\n:8080", BLEU_FONCE),
        ("Cassandra\n:9042", BLEU_MED),
        ("RabbitMQ\n:5672", VERT),
        ("MinIO\n:9002", ORANGE),
        ("Redis\n:6379", ROUGE),
        ("Prometheus\n:9090", BLEU_MED),
        ("Grafana\n:3001", ORANGE),
    ]
    add_text_box(sl, "INFRASTRUCTURE", Inches(0.5), Inches(5.1), Inches(12), Inches(0.35),
                 size=10, bold=True, color=GRIS)
    add_rect(sl, Inches(0.3), Inches(5.1), W - Inches(0.5), Inches(0.04), GRIS)
    for i, (name, col) in enumerate(infra_svcs):
        x = Inches(0.35) + i * Inches(1.85)
        y = Inches(5.5)
        add_rect(sl, x, y, Inches(1.7), Inches(0.75), col)
        add_text_box(sl, name, x, y, Inches(1.7), Inches(0.75),
                     size=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    # RabbitMQ arrow
    add_text_box(sl, "Événements AMQP (asynchrone)", Inches(3.5), Inches(4.7), Inches(4.0), Inches(0.35),
                 size=9, italic=True, color=VERT)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 6: INFRASTRUCTURE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, VERT)
    add_left_accent(sl, VERT)

    add_text_box(sl, "03  Infrastructure Technique", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), VERT)

    infra_cards = [
        ("Keycloak 26", "IAM / OpenID Connect\nRealm ent-est · JWKS RS256\nRôles : admin, teacher, student", BLEU_FONCE),
        ("PostgreSQL 16", "Backend relationnel\nde Keycloak\nPort 5432", BLEU_MED),
        ("Cassandra 5.0", "Base NoSQL large-colonne\nKeyspace ent_est\n11 tables applicatives", VERT),
        ("RabbitMQ 4", "Broker AMQP\nExchange topic ent.events.topic\nEvents versionnés (.v1)", ORANGE),
        ("MinIO", "Stockage objet S3\n3 buckets : courses, uploads, logs\nPresigned URLs TTL 180s", VERT),
        ("Redis 7", "Cache distribué\nRate limiting\nPort 6379", BLEU_MED),
        ("Mailpit v1.18", "SMTP sandbox dev\nCapture emails sans envoi réel\nUI port 8025", ORANGE),
        ("Prometheus", "Scrape /metrics 15s\n9 services instrumentés\nPort 9090", ROUGE),
        ("Grafana", "Dashboards temps-réel\nDatasource auto-provisionnée\nPort 3001", BLEU_MED),
    ]

    for i, (name, desc, col) in enumerate(infra_cards):
        col_i = i % 3; row_i = i // 3
        x = Inches(0.3) + col_i * Inches(4.35)
        y = Inches(1.0) + row_i * Inches(1.95)
        add_rect(sl, x, y, Inches(4.15), Inches(1.8), BLANC)
        add_rect(sl, x, y, Inches(0.06), Inches(1.8), col)
        add_rect(sl, x, y, Inches(4.15), Inches(0.45), col)
        add_text_box(sl, name, x + Inches(0.12), y + Inches(0.07), Inches(3.9), Inches(0.35),
                     size=13, bold=True, color=BLANC)
        add_text_box(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(3.85), Inches(1.1),
                     size=10, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 7: MICROSERVICES 1
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, ORANGE)
    add_left_accent(sl, ORANGE)

    add_text_box(sl, "04  Microservices — Services Métier", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), ORANGE)

    services_detail = [
        ("ms-auth-core", "8010", "Interne", "Validation JWT via JWKS Keycloak\n/auth/me → claims complètes", BLEU_FONCE),
        ("ms-identity-admin", "8013", "admin", "CRUD utilisateurs Keycloak + Cassandra\nPublie user.created.v1", BLEU_MED),
        ("ms-course-content", "8011", "teacher, admin", "Gestion cours + assets MinIO\nPublie course.created.v1 + asset.uploaded.v1", VERT),
        ("ms-course-access", "8012", "student", "Lecture cours + presigned URLs\nTTL 180s, délègue à course-content", VERT),
        ("ms-notification", "8014", "Consumer AMQP", "RabbitMQ consumer + SMTP Mailpit\nStocke notifications en Cassandra", ORANGE),
        ("ms-calendar-schedule", "8015", "teacher/admin (✍), tous (📖)", "CRUD événements calendrier\nFiltres ?month + ?module_code", BLEU_MED),
        ("ms-forum-chat", "8016", "Tous", "Forum REST + WebSocket chat\nSalles par module_code", ORANGE),
        ("ms-exam-assignment", "8017", "teacher/student", "Devoirs, soumissions MinIO, notes\nPublie assignment.created + grade.assigned", ROUGE),
        ("ms-ai-assistant", "8018", "Tous", "Chat IA / Résumé / FAQ\nOllama Llama3, rate limit 20 req/min", BLEU_FONCE),
    ]

    for i, (name, port, roles, desc, col) in enumerate(services_detail):
        col_i = i % 3; row_i = i // 3
        x = Inches(0.3) + col_i * Inches(4.35)
        y = Inches(1.0) + row_i * Inches(1.95)
        add_rect(sl, x, y, Inches(4.15), Inches(1.8), BLANC)
        add_rect(sl, x, y, Inches(0.06), Inches(1.8), col)
        # header
        add_rect(sl, x + Inches(0.06), y, Inches(4.09), Inches(0.5), RGBColor(0xF8,0xFA,0xFF))
        add_text_box(sl, name, x + Inches(0.12), y + Inches(0.05), Inches(2.5), Inches(0.38),
                     size=11, bold=True, color=col)
        add_rect(sl, x + Inches(2.7), y + Inches(0.08), Inches(0.7), Inches(0.32), col)
        add_text_box(sl, f":{port}", x + Inches(2.7), y + Inches(0.08), Inches(0.7), Inches(0.32),
                     size=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        add_text_box(sl, roles, x + Inches(0.12), y + Inches(0.52), Inches(3.9), Inches(0.28),
                     size=8.5, italic=True, color=col)
        add_text_box(sl, desc, x + Inches(0.12), y + Inches(0.85), Inches(3.9), Inches(0.85),
                     size=9.5, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 8: GATEWAY & SECURITE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_MED)
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "05  API Gateway & Sécurité", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), BLEU_MED)

    # Left: Gateway flow
    card(sl, Inches(0.3), Inches(1.0), Inches(6.1), Inches(5.8), bg=BLANC, border_color=BLEU_FONCE,
         title="Flux d'Authentification JWT", title_size=13)

    steps = [
        ("1", "Frontend envoie Authorization: Bearer <JWT>", BLEU_FONCE),
        ("2", "Gateway appelle GET /auth/me sur ms-auth-core", BLEU_MED),
        ("3", "ms-auth-core valide signature via JWKS Keycloak", VERT),
        ("4", "Claims retournées : sub, username, realm_access.roles", VERT),
        ("5", "require_roles() vérifie les permissions", ORANGE),
        ("6", "Si autorisé → proxy vers le microservice cible", VERT),
        ("7", "Sinon → 401 Unauthorized ou 403 Forbidden", ROUGE),
    ]
    for i, (num, step, col) in enumerate(steps):
        y_pos = Inches(1.55 + i * 0.68)
        add_rect(sl, Inches(0.5), y_pos, Inches(0.32), Inches(0.32), col)
        add_text_box(sl, num, Inches(0.5), y_pos, Inches(0.32), Inches(0.32),
                     size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        add_text_box(sl, step, Inches(0.9), y_pos, Inches(5.3), Inches(0.32),
                     size=10.5, color=GRIS)

    # Right: RBAC + resilience
    card(sl, Inches(6.7), Inches(1.0), Inches(6.3), Inches(2.6), bg=BLANC, border_color=BLEU_MED,
         title="RBAC — Rôles et Permissions", title_size=13)
    rbac_data = [
        ("admin", "Tous les accès : utilisateurs, cours, calendrier", BLEU_FONCE),
        ("teacher", "Création cours/devoirs, calendrier, forum", VERT),
        ("student", "Lecture cours, soumissions, forum, calendrier", ORANGE),
    ]
    for i, (role, perm, col) in enumerate(rbac_data):
        y_pos = Inches(1.55 + i * 0.7)
        add_rect(sl, Inches(6.85), y_pos, Inches(1.0), Inches(0.3), col)
        add_text_box(sl, role, Inches(6.85), y_pos, Inches(1.0), Inches(0.3),
                     size=10, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        add_text_box(sl, perm, Inches(8.0), y_pos, Inches(4.8), Inches(0.3),
                     size=9.5, color=GRIS)

    card(sl, Inches(6.7), Inches(3.8), Inches(6.3), Inches(2.8), bg=BLANC, border_color=ORANGE,
         title="Résilience & Sécurité", title_size=13)
    resilience = [
        "Retry 2x sur GET avec backoff exponentiel (150ms→300ms)",
        "Timeout configurable par service (défaut 8s)",
        "CORS strict : 3 origines autorisées seulement",
        "Réseau Docker interne (services non exposés)",
        "INTERNAL_API_TOKEN pour appels inter-services",
        "Variables sensibles dans .env.compose (non versionné)",
    ]
    for i, item in enumerate(resilience):
        add_text_box(sl, f"  • {item}", Inches(6.85), Inches(4.35 + i * 0.38), Inches(6.0), Inches(0.35),
                     size=10, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 9: FRONTEND
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, VERT)
    add_left_accent(sl, VERT)

    add_text_box(sl, "06  Interface Utilisateur — Frontend", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), VERT)

    # Stack badges
    stack_badges = [
        ("React 18", BLEU_FONCE), ("TypeScript", BLEU_MED), ("Vite", ORANGE),
        ("Tailwind CSS", VERT), ("shadcn/ui", BLEU_MED), ("React Router v6", VERT),
        ("WebSocket", ROUGE), ("nginx", GRIS),
    ]
    for i, (tech, col) in enumerate(stack_badges):
        x = Inches(0.35) + (i % 4) * Inches(3.3)
        y = Inches(1.0) + (i // 4) * Inches(0.5)
        badge(sl, tech, x, y, Inches(2.6), Inches(0.38), bg=col, size=11)

    # Pages grid
    add_text_box(sl, "Pages de l'Application", Inches(0.3), Inches(2.15), Inches(12), Inches(0.4),
                 size=14, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(2.55), Inches(1.5), Inches(0.04), VERT)

    pages = [
        ("Dashboard", "Stats temps-réel, fil activité", BLEU_FONCE),
        ("Admin", "Gestion utilisateurs & rôles", BLEU_MED),
        ("Teacher", "Cours, assets, devoirs", VERT),
        ("Student", "Cours, téléchargements", VERT),
        ("Calendar", "Calendrier mensuel interactif", ORANGE),
        ("Forum", "Discussions + Chat WebSocket", ORANGE),
        ("Exams", "Devoirs, soumissions, notes", ROUGE),
        ("Notifications", "Centre de notifications", BLEU_MED),
        ("Profile", "Profil & déconnexion", GRIS),
        ("AI Assistant", "Chat IA, résumé, FAQ", BLEU_FONCE),
    ]
    for i, (name, desc, col) in enumerate(pages):
        col_i = i % 5; row_i = i // 5
        x = Inches(0.3) + col_i * Inches(2.62)
        y = Inches(2.75) + row_i * Inches(1.75)
        add_rect(sl, x, y, Inches(2.5), Inches(1.6), BLANC)
        add_rect(sl, x, y, Inches(2.5), Inches(0.45), col)
        add_text_box(sl, name, x + Inches(0.08), y + Inches(0.06), Inches(2.3), Inches(0.35),
                     size=11, bold=True, color=BLANC)
        add_text_box(sl, desc, x + Inches(0.08), y + Inches(0.55), Inches(2.3), Inches(0.9),
                     size=10, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 10: RABBITMQ EVENTS
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, ORANGE)
    add_left_accent(sl, ORANGE)

    add_text_box(sl, "04  Messagerie Asynchrone — RabbitMQ Events", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=22, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(2.0), Inches(0.04), ORANGE)

    # Left: Envelope schema
    card(sl, Inches(0.3), Inches(1.0), Inches(5.5), Inches(4.0), bg=BLANC, border_color=ORANGE,
         title="Enveloppe Standard d'Événement", title_size=13)
    code = (
        '{\n'
        '  "event_id":       "<uuid>",\n'
        '  "event_type":     "calendar.event.v1",\n'
        '  "occurred_at":    "2026-04-24T10:00:00Z",\n'
        '  "producer":       "ms-calendar-schedule",\n'
        '  "correlation_id": "<uuid>",\n'
        '  "payload":        { ... }\n'
        '}'
    )
    add_rect(sl, Inches(0.5), Inches(1.55), Inches(5.0), Inches(2.8), RGBColor(0x1E,0x29,0x3B))
    add_text_box(sl, code, Inches(0.6), Inches(1.65), Inches(4.8), Inches(2.6),
                 size=10, color=BLANC)

    # Right: Events list
    card(sl, Inches(6.1), Inches(1.0), Inches(6.9), Inches(5.8), bg=BLANC, border_color=ORANGE,
         title="Événements Publiés par Service", title_size=13)

    events = [
        ("ms-identity-admin",     "user.created.v1",            BLEU_FONCE),
        ("ms-course-content",     "course.created.v1",          BLEU_MED),
        ("ms-course-content",     "asset.uploaded.v1",          BLEU_MED),
        ("ms-calendar-schedule",  "calendar.event.created.v1",  VERT),
        ("ms-forum-chat",         "forum.thread.created.v1",    ORANGE),
        ("ms-exam-assignment",    "assignment.created.v1",      ROUGE),
        ("ms-exam-assignment",    "grade.assigned.v1",          ROUGE),
    ]
    for i, (producer, event, col) in enumerate(events):
        y_pos = Inches(1.55 + i * 0.65)
        add_rect(sl, Inches(6.25), y_pos, Inches(0.3), Inches(0.3), col)
        add_text_box(sl, producer, Inches(6.65), y_pos, Inches(2.5), Inches(0.3),
                     size=8.5, italic=True, color=GRIS)
        add_text_box(sl, f"→ {event}", Inches(9.1), y_pos, Inches(3.7), Inches(0.3),
                     size=9, bold=True, color=col)

    # Consumer
    add_text_box(sl, "Consommateur : ms-notification", Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.35),
                 size=11, bold=True, color=VERT)
    add_text_box(sl, "Queue: q.notification.user  ·  Exchange: ent.events.topic  ·  Pattern: #",
                 Inches(0.5), Inches(5.55), Inches(5.5), Inches(0.35), size=10, color=GRIS)
    add_text_box(sl, "→ Email SMTP (Mailpit)  +  Stockage Cassandra (table notifications)",
                 Inches(0.5), Inches(5.9), Inches(5.5), Inches(0.35), size=10, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 11: OBSERVABILITE
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, ORANGE)
    add_left_accent(sl, ORANGE)

    add_text_box(sl, "07  Observabilité — Prometheus & Grafana", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=22, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(2.0), Inches(0.04), ORANGE)

    # Prometheus targets
    card(sl, Inches(0.3), Inches(1.0), Inches(5.8), Inches(5.8), bg=BLANC, border_color=ORANGE,
         title="Prometheus — Scrape Configuration", title_size=13)
    targets = [
        "ms-auth-core:8000/metrics",
        "ms-api-gateway:8000/metrics",
        "ms-identity-admin:8000/metrics",
        "ms-course-content:8000/metrics",
        "ms-course-access:8000/metrics",
        "ms-notification:8000/metrics",
        "ms-calendar-schedule:8000/metrics",
        "ms-forum-chat:8000/metrics",
        "ms-exam-assignment:8000/metrics",
    ]
    add_text_box(sl, "scrape_interval: 15s  |  evaluation_interval: 15s",
                 Inches(0.5), Inches(1.5), Inches(5.3), Inches(0.35), size=9, italic=True, color=GRIS)
    for i, t in enumerate(targets):
        add_rect(sl, Inches(0.5), Inches(1.95 + i * 0.52), Inches(0.28), Inches(0.28), ORANGE)
        add_text_box(sl, t, Inches(0.88), Inches(1.92 + i * 0.52), Inches(5.0), Inches(0.3),
                     size=10, color=GRIS)

    # Metrics collected
    card(sl, Inches(6.4), Inches(1.0), Inches(6.6), Inches(2.7), bg=BLANC, border_color=BLEU_MED,
         title="Métriques Collectées (runtime.py)", title_size=13)
    metrics = [
        "Compteur requêtes HTTP (méthode, route, status)",
        "Histogramme latence (p50, p95, p99)",
        "Compteur erreurs par type d'exception",
        "Métriques système (CPU, mémoire)",
    ]
    for i, m in enumerate(metrics):
        add_text_box(sl, f"  • {m}", Inches(6.55), Inches(1.55 + i * 0.52), Inches(6.2), Inches(0.4),
                     size=10.5, color=GRIS)

    # Grafana
    card(sl, Inches(6.4), Inches(3.9), Inches(6.6), Inches(2.8), bg=BLANC, border_color=ORANGE,
         title="Grafana — Dashboards", title_size=13)
    grafana_items = [
        "Port 3001, admin/ChangeMe_123!",
        "Datasource Prometheus auto-provisionnée",
        "Dashboard ENT Platform (ent-platform.json)",
        "Provisionnement via /etc/grafana/provisioning",
        "Alertes configurables sur les métriques critiques",
    ]
    for i, g in enumerate(grafana_items):
        add_text_box(sl, f"  • {g}", Inches(6.55), Inches(4.45 + i * 0.48), Inches(6.2), Inches(0.38),
                     size=10.5, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 12: TESTS & DEPLOIEMENT
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, ROUGE)
    add_left_accent(sl, ROUGE)

    add_text_box(sl, "08  Tests & Déploiement", Inches(0.3), Inches(0.2), Inches(12), Inches(0.55),
                 size=24, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(1.5), Inches(0.04), ROUGE)

    # Tests
    card(sl, Inches(0.3), Inches(1.0), Inches(6.3), Inches(5.8), bg=BLANC, border_color=VERT,
         title="Scripts de Smoke-Tests (scripts/)", title_size=13)
    scripts_list = [
        ("compose-healthcheck.sh", "Vérification tous les conteneurs"),
        ("rbac-smoke-test.sh", "RBAC admin/teacher/student"),
        ("course-content-smoke-test.sh", "CRUD cours + upload asset"),
        ("course-access-smoke-test.sh", "Listing + presigned URLs"),
        ("identity-admin-smoke-test.sh", "Création user + rôles"),
        ("calendar-smoke-test.sh", "CRUD événements calendrier"),
        ("forum-smoke-test.sh", "Forum + WebSocket chat"),
        ("exam-smoke-test.sh", "Devoirs + soumissions + notes"),
    ]
    for i, (script, desc) in enumerate(scripts_list):
        add_rect(sl, Inches(0.5), Inches(1.55 + i * 0.62), Inches(0.28), Inches(0.28), VERT)
        add_text_box(sl, script, Inches(0.88), Inches(1.52 + i * 0.62), Inches(2.8), Inches(0.3),
                     size=9.5, bold=True, color=VERT)
        add_text_box(sl, f"→ {desc}", Inches(3.75), Inches(1.52 + i * 0.62), Inches(2.6), Inches(0.3),
                     size=9.5, color=GRIS)

    # Deploiement
    card(sl, Inches(6.85), Inches(1.0), Inches(6.15), Inches(2.9), bg=BLANC, border_color=BLEU_MED,
         title="Commandes de Déploiement", title_size=13)
    deploy_cmds = [
        "# Démarrer tous les services",
        "docker compose --env-file .env.compose \\",
        "  up -d --build",
        "",
        "# Vérifier l'état",
        "docker compose ps",
        "",
        "# Logs en temps-réel",
        "docker compose logs -f ms-api-gateway",
    ]
    add_rect(sl, Inches(7.0), Inches(1.55), Inches(5.8), Inches(2.1), RGBColor(0x1E,0x29,0x3B))
    add_text_box(sl, "\n".join(deploy_cmds), Inches(7.1), Inches(1.6), Inches(5.6), Inches(2.0),
                 size=9, color=BLANC)

    card(sl, Inches(6.85), Inches(4.1), Inches(6.15), Inches(2.7), bg=BLANC, border_color=ORANGE,
         title="Ordre de Démarrage Docker", title_size=13)
    order = [
        "1. postgres → keycloak → bootstrap scripts",
        "2. cassandra → cassandra-init",
        "3. minio → minio-init",
        "4. rabbitmq",
        "5. ms-auth-core (healthy) → microservices",
        "6. microservices → ms-api-gateway",
        "7. ms-api-gateway → frontend",
    ]
    for i, step in enumerate(order):
        add_text_box(sl, step, Inches(7.0), Inches(4.65 + i * 0.33), Inches(5.9), Inches(0.3),
                     size=9.5, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 13: PATTERNS
    sl = blank_slide(prs)
    slide_bg(sl, FOND)
    add_top_bar(sl, BLEU_FONCE)
    add_left_accent(sl, BLEU_MED)

    add_text_box(sl, "Patterns & Conventions — Structure Uniforme des Services", Inches(0.3), Inches(0.2), Inches(12.5), Inches(0.55),
                 size=20, bold=True, color=BLEU_FONCE)
    add_rect(sl, Inches(0.3), Inches(0.8), Inches(2.5), Inches(0.04), BLEU_MED)

    # File structure
    card(sl, Inches(0.3), Inches(1.0), Inches(4.0), Inches(5.8), bg=BLANC, border_color=BLEU_FONCE,
         title="Structure de chaque service", title_size=12)
    structure = [
        "services/<name>/",
        "  ├── app/",
        "  │   ├── __init__.py",
        "  │   ├── runtime.py  ← copie identique",
        "  │   └── main.py     ← logique métier",
        "  ├── requirements.txt",
        "  └── Dockerfile",
        "",
        "gateway/",
        "  ├── app/",
        "  │   ├── main.py     ← toutes les routes",
        "  │   └── runtime.py",
        "  └── Dockerfile",
    ]
    add_rect(sl, Inches(0.5), Inches(1.55), Inches(3.6), Inches(4.0), RGBColor(0x1E,0x29,0x3B))
    add_text_box(sl, "\n".join(structure), Inches(0.55), Inches(1.6), Inches(3.5), Inches(3.9),
                 size=9, color=BLANC)

    # Patterns
    patterns = [
        ("runtime.py Partagé", BLEU_MED,
         "• Logging structuré JSON + request_id\n• Middleware Prometheus automatique\n• GET /metrics exposition\n• Gestionnaire exceptions global"),
        ("Cassandra _ensure_db()", VERT,
         "• Init paresseuse au premier appel API\n• CREATE KEYSPACE IF NOT EXISTS ent_est\n• CREATE TABLE IF NOT EXISTS\n• PlainTextAuthProvider"),
        ("RabbitMQ Event Envelope", ORANGE,
         "• event_id, event_type, occurred_at\n• producer, correlation_id, payload\n• Versionnement .v1\n• Exchange topic durable"),
        ("Auth Pattern", ROUGE,
         "• GET /auth/me avec Bearer token\n• Check realm_access.roles\n• require_roles() comme FastAPI Depends\n• JWT propagé vers services"),
    ]
    for i, (title, col, desc) in enumerate(patterns):
        col_i = i % 2; row_i = i // 2
        x = Inches(4.5) + col_i * Inches(4.45)
        y = Inches(1.0) + row_i * Inches(2.85)
        add_rect(sl, x, y, Inches(4.3), Inches(2.7), BLANC)
        add_rect(sl, x, y, Inches(4.3), Inches(0.45), col)
        add_rect(sl, x, y, Inches(0.06), Inches(2.7), col)
        add_text_box(sl, title, x + Inches(0.12), y + Inches(0.06), Inches(4.0), Inches(0.35),
                     size=12, bold=True, color=BLANC)
        add_text_box(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(4.0), Inches(2.0),
                     size=10, color=GRIS)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), BLEU_FONCE)
    add_text_box(sl, "ENT-EST-Salé  ·  EST Salé", 0, H - Inches(0.28), W, Inches(0.25),
                 size=9, color=BLANC, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────── SLIDE 14: CONCLUSION
    sl = blank_slide(prs)
    slide_bg(sl, BLEU_FONCE)

    add_rect(sl, W * 0.65, 0, W * 0.35, H, BLEU_MED)
    add_rect(sl, W * 0.63, 0, Inches(0.04), H, BLANC)

    add_text_box(sl, "09", Inches(0.5), Inches(0.6), Inches(1.2), Inches(0.9),
                 size=52, bold=True, color=BLEU_MED)
    add_text_box(sl, "Conclusion &\nPerspectives", Inches(0.5), Inches(1.4), Inches(7.5), Inches(1.2),
                 size=32, bold=True, color=BLANC)
    add_rect(sl, Inches(0.5), Inches(2.6), Inches(3.0), Inches(0.05), BLEU_MED)

    realisations = [
        "10 microservices Python FastAPI entièrement opérationnels",
        "Sécurité OAuth2/OIDC avec Keycloak + RBAC 3 rôles",
        "Persistance NoSQL Apache Cassandra (11 tables)",
        "Messagerie asynchrone RabbitMQ (7 types d'événements)",
        "Stockage objet MinIO S3 (3 buckets, presigned URLs)",
        "Frontend React 18 + shadcn/ui (11 pages rôle-adaptatives)",
        "Chat temps-réel WebSocket par salle de module",
        "Observabilité Prometheus + Grafana (9 services)",
        "8 scripts de smoke-test couvrant tous les flux",
        "Déploiement one-command Docker Compose",
    ]
    add_text_box(sl, "Réalisations", Inches(0.5), Inches(2.8), Inches(7.5), Inches(0.4),
                 size=13, bold=True, color=BLEU_CLAIR)
    for i, real in enumerate(realisations):
        add_rect(sl, Inches(0.5), Inches(3.25 + i * 0.38), Inches(0.22), Inches(0.22), VERT)
        add_text_box(sl, real, Inches(0.82), Inches(3.22 + i * 0.38), Inches(7.0), Inches(0.3),
                     size=10.5, color=BLANC)

    # Right: Perspectives
    add_text_box(sl, "Perspectives", W * 0.67, Inches(0.5), Inches(4.5), Inches(0.45),
                 size=14, bold=True, color=BLANC)
    perspectives = [
        "Migration Kubernetes",
        "(Helm charts préparés)",
        "CI/CD GitHub Actions",
        "Dashboards Grafana avancés",
        "App mobile React Native",
        "Cassandra multi-nœuds HA",
        "Assistant IA en production",
    ]
    for i, p in enumerate(perspectives):
        add_text_box(sl, f"  → {p}", W * 0.67, Inches(1.0 + i * 0.55), Inches(4.5), Inches(0.42),
                     size=11, color=BLANC)

    add_text_box(sl, "Merci pour votre attention", Inches(0.5), Inches(6.5), Inches(7.5), Inches(0.7),
                 size=20, bold=True, italic=True, color=BLEU_CLAIR)

    add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), RGBColor(0x09, 0x1B, 0x36))
    add_text_box(sl, "ENT-EST-Salé  ·  DevOps et Cloud  ·  EST Salé — Université Mohammed V  ·  Avril 2026",
                 0, H - Inches(0.28), W, Inches(0.25), size=9, color=BLEU_CLAIR, align=PP_ALIGN.CENTER)

    return prs


# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = build_presentation()
    out = "/home/khalidmarzoug/Documents/devops-project/ent-est-sale-platform/Presentation_ENT_EST_Sale.pptx"
    prs.save(out)
    print(f"✅  Présentation sauvegardée : {out}")
